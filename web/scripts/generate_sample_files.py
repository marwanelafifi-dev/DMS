from __future__ import annotations

from datetime import datetime, timezone
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile, ZipInfo

from docx import Document
from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


OUTPUT_DIR = Path(__file__).resolve().parents[1] / "public" / "sample-files"
FIXED_ZIP_TIME = (2026, 7, 27, 8, 0, 0)
FIXED_DOCUMENT_TIME = datetime(2026, 7, 27, 8, 0, tzinfo=timezone.utc)


def normalize_office_archive(path: Path) -> None:
    """Make generated Office archives reproducible across local runs."""
    temporary_path = path.with_suffix(f"{path.suffix}.tmp")
    with ZipFile(path, "r") as source, ZipFile(temporary_path, "w", ZIP_DEFLATED) as target:
        for source_info in sorted(source.infolist(), key=lambda item: item.filename):
            target_info = ZipInfo(source_info.filename, FIXED_ZIP_TIME)
            target_info.compress_type = ZIP_DEFLATED
            target_info.external_attr = source_info.external_attr
            target_info.create_system = source_info.create_system
            target.writestr(target_info, source.read(source_info.filename))
    temporary_path.replace(path)


def create_text_sample() -> None:
    content = """DMS - Enterprise Document Management System
============================================

Sample Text File for Testing

This is a sample text document used to test the preview functionality
of the Enterprise Document Management System (DMS).

The system supports the following file types:
- Text Files (.txt)
- Word Documents (.docx, .doc)
- Excel Spreadsheets (.xlsx, .xls)
- PowerPoint Presentations (.pptx, .ppt)
- PDF Documents (.pdf)
- Images (.png, .jpg, .jpeg, .gif)

Quality Management System (QMS)
- ISO 9001:2015 Compliance
- Document Control
- Change Management
- Version History

Information Security Management (ISMS)
- ISO 27001:2022 Compliance
- Access Control
- Audit Logging
- Risk Assessment
"""
    (OUTPUT_DIR / "DMS-Sample-Text.txt").write_text(content, encoding="utf-8")


def create_word_sample() -> None:
    path = OUTPUT_DIR / "DMS-Sample-Document.docx"
    document = Document()
    document.core_properties.title = "DMS Sample Controlled Document"
    document.core_properties.subject = "Local document parsing and preview verification"
    document.core_properties.author = "DMS Sample Pack"
    document.core_properties.created = FIXED_DOCUMENT_TIME
    document.core_properties.modified = FIXED_DOCUMENT_TIME
    document.add_heading("DMS Sample Controlled Document", level=0)
    document.add_paragraph("Tracking code: DOCX-SAMPLE-4729")
    document.add_heading("Purpose", level=1)
    document.add_paragraph(
        "Verify local Docling conversion, Markdown rendering, persistent preview, "
        "download, metadata search, and approval workflows."
    )
    document.add_heading("Review checklist", level=1)
    for item in (
        "Upload completes without leaving the Document Library.",
        "Extracted content appears in the read-only preview.",
        "The document remains available after refreshing the browser.",
    ):
        document.add_paragraph(item, style="List Bullet")
    document.save(path)
    normalize_office_archive(path)


def create_excel_sample() -> None:
    path = OUTPUT_DIR / "DMS-Sample-Spreadsheet.xlsx"
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Quality Metrics"
    worksheet.append(["Metric", "Target", "Actual", "Status"])
    worksheet.append(["First-pass yield", "98.0%", "98.7%", "On target"])
    worksheet.append(["CAPA closure", "30 days", "27 days", "On target"])
    worksheet.append(["Training compliance", "100%", "99.2%", "Review"])
    worksheet.append(["Tracking code", "XLSX-SAMPLE-4729", "", ""])
    for cell in worksheet[1]:
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="2F3E83")
        cell.alignment = Alignment(horizontal="center")
    worksheet.freeze_panes = "A2"
    worksheet.column_dimensions["A"].width = 24
    worksheet.column_dimensions["B"].width = 20
    worksheet.column_dimensions["C"].width = 16
    worksheet.column_dimensions["D"].width = 16
    workbook.properties.title = "DMS Sample Quality Metrics"
    workbook.properties.creator = "DMS Sample Pack"
    workbook.properties.created = FIXED_DOCUMENT_TIME.replace(tzinfo=None)
    workbook.properties.modified = FIXED_DOCUMENT_TIME.replace(tzinfo=None)
    workbook.save(path)
    normalize_office_archive(path)


def create_powerpoint_sample() -> None:
    path = OUTPUT_DIR / "DMS-Sample-Presentation.pptx"
    presentation = Presentation()
    presentation.core_properties.title = "DMS Local Preview Sample"
    presentation.core_properties.subject = "PowerPoint parsing verification"
    presentation.core_properties.author = "DMS Sample Pack"
    presentation.core_properties.created = FIXED_DOCUMENT_TIME
    presentation.core_properties.modified = FIXED_DOCUMENT_TIME

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = "DMS Local Preview Sample"
    title_slide.placeholders[1].text = "PPTX-SAMPLE-4729 · Docling and workflow verification"

    content_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    content_slide.shapes.title.text = "What this sample verifies"
    text_frame = content_slide.placeholders[1].text_frame
    text_frame.clear()
    for index, text in enumerate(
        (
            "PowerPoint upload and local Docling conversion",
            "Formatted Markdown preview in the Document Library",
            "Persistent source retrieval after browser refresh",
            "OCR content search and read-only download",
        )
    ):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = text
        paragraph.level = 0

    accent = content_slide.shapes.add_shape(
        1,
        Inches(0.55),
        Inches(6.6),
        Inches(12.2),
        Inches(0.35),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(63, 139, 202)
    accent.line.fill.background()
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Arial"
                    if shape == slide.shapes.title:
                        run.font.size = Pt(28)

    presentation.save(path)
    normalize_office_archive(path)


def create_pdf_sample() -> None:
    path = OUTPUT_DIR / "DMS-Sample-Report.pdf"
    lines = [
        "DMS Local PDF Sample",
        "Tracking code: PDF-SAMPLE-4729",
        "This controlled report verifies browser PDF preview,",
        "local Docling extraction, OCR search, and persistent download.",
    ]
    stream_parts = ["BT", "/F1 22 Tf", "72 720 Td", f"({lines[0]}) Tj", "/F1 12 Tf"]
    for line in lines[1:]:
        escaped = line.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
        stream_parts.extend(["0 -28 Td", f"({escaped}) Tj"])
    stream_parts.append("ET")
    stream = "\n".join(stream_parts).encode("ascii")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>"
        ),
        b"<< /Length " + str(len(stream)).encode("ascii") + b" >>\nstream\n" + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    pdf = bytearray(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets = [0]
    for index, body in enumerate(objects, start=1):
        offsets.append(len(pdf))
        pdf.extend(f"{index} 0 obj\n".encode("ascii"))
        pdf.extend(body)
        pdf.extend(b"\nendobj\n")
    xref_offset = len(pdf)
    pdf.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    pdf.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        pdf.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    pdf.extend(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode("ascii")
    )
    path.write_bytes(pdf)


def sample_font(size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    for font_path in (
        "C:/Windows/Fonts/arialbd.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
    ):
        if Path(font_path).exists():
            return ImageFont.truetype(font_path, size=size)
    return ImageFont.load_default()


def create_image_samples() -> None:
    width, height = 1280, 720
    image = Image.new("RGB", (width, height), "#0f172a")
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle((70, 70, 1210, 650), radius=38, fill="#f8fafc")
    draw.rectangle((70, 70, 1210, 185), fill="#2f3e83")
    draw.text((115, 100), "DMS SAMPLE IMAGE", fill="white", font=sample_font(48))
    draw.text((115, 265), "LOCAL OCR SAMPLE 4729", fill="#1e293b", font=sample_font(54))
    draw.text((115, 365), "PNG preview and text extraction", fill="#3f8bca", font=sample_font(34))
    draw.text((115, 435), "Refresh the browser and open this image again.", fill="#475569", font=sample_font(28))
    image.save(OUTPUT_DIR / "DMS-Sample-Image.png", format="PNG", optimize=True)

    photo = Image.new("RGB", (width, height))
    pixels = photo.load()
    for y in range(height):
        for x in range(width):
            pixels[x, y] = (
                35 + (x * 35 // width),
                62 + (y * 70 // height),
                131 + ((x + y) * 70 // (width + height)),
            )
    photo_draw = ImageDraw.Draw(photo)
    photo_draw.rounded_rectangle((110, 145, 1170, 575), radius=44, fill="#ffffff")
    photo_draw.text((180, 225), "DMS JPG SAMPLE", fill="#2f3e83", font=sample_font(56))
    photo_draw.text((180, 335), "PHOTO-OCR-4729", fill="#0f172a", font=sample_font(44))
    photo_draw.text((180, 425), "Image preview verification", fill="#3f8bca", font=sample_font(30))
    photo.save(OUTPUT_DIR / "DMS-Sample-Photo.jpg", format="JPEG", quality=90, optimize=True)


def create_readme() -> None:
    content = """# DMS Sample Files

These files exercise the real local upload, Docling conversion, OCR search, preview,
download, persistence, and workflow paths.

- `DMS-Sample-Text.txt` — plain text
- `DMS-Sample-Document.docx` — Microsoft Word
- `DMS-Sample-Spreadsheet.xlsx` — Microsoft Excel
- `DMS-Sample-Presentation.pptx` — Microsoft PowerPoint
- `DMS-Sample-Report.pdf` — browser PDF preview
- `DMS-Sample-Image.png` — PNG image and OCR text
- `DMS-Sample-Photo.jpg` — JPEG image and OCR text

Regenerate the pack from `web/` with:

```powershell
python scripts/generate_sample_files.py
```
"""
    (OUTPUT_DIR / "README.md").write_text(content, encoding="utf-8")


def main() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    create_text_sample()
    create_word_sample()
    create_excel_sample()
    create_powerpoint_sample()
    create_pdf_sample()
    create_image_samples()
    create_readme()
    print(f"Generated DMS sample pack in {OUTPUT_DIR}")


if __name__ == "__main__":
    main()
